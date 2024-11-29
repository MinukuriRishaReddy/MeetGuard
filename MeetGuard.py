from pptx import Presentation

# Create a new PowerPoint presentation
presentation = Presentation()

# Slide 1: Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Real-Time Sign Language Interpreter"
subtitle.text = "Bridging Communication Gaps in Virtual Meetings"

# Slide 2: Introduction
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = (
    "• Overview of the problem\n"
    "• Purpose of the tool\n"
    "• Objectives"
)

# Slide 3: Key Features (1)
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Key Features (1)"
content.text = (
    "• Real-Time Gesture Recognition\n"
    "• Text and Speech Output"
)

# Slide 4: Key Features (2)

